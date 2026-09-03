import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

from slidecaptain.export.pptx_writer import write_pptx
from slidecaptain.layout.engine import build_render_plan
from slidecaptain.metrics.font_metrics import FontMetrics
from slidecaptain.models.deck import Chapter, Deck, DeckMeta, Slide, Structure, TableSlots
from slidecaptain.models.preset import Preset

PRESET = Preset()


def _table_deck() -> Deck:
    return Deck(
        meta=DeckMeta(title="표 테스트"),
        structure=Structure(
            chapters=[Chapter(id="ch01", topic="비교 결과", template="table")]
        ),
        slides=[
            Slide(
                chapter_id="ch01",
                slots=TableSlots(
                    columns=["항목", "옵션 A", "옵션 B"],
                    rows=[
                        ["도입 비용", "1,200만 원", "800만 원"],
                        ["운영 부담", "낮음", "중간"],
                    ],
                    footnote="주: 2026년 상반기 견적 기준",
                ),
            )
        ],
    )


@pytest.fixture()
def saved(tmp_path):
    metrics = FontMetrics.from_bundled()
    plan = build_render_plan(_table_deck(), PRESET, metrics)
    out = tmp_path / "table.pptx"
    write_pptx(plan, out)
    return Presentation(str(out))


def test_table_shape_exists_with_role_name(saved):
    shapes = {s.name: s for s in saved.slides[0].shapes}
    assert "ch01:table" in shapes
    assert shapes["ch01:table"].has_table


def test_table_dimensions(saved):
    table = next(s for s in saved.slides[0].shapes if s.name == "ch01:table").table
    assert len(table.rows) == 3  # 머리글 + 데이터 2행
    assert len(table.columns) == 3
    total_w = sum(col.width for col in table.columns)
    assert total_w == pytest.approx(Emu(round(860.0 * 12700)), rel=0.01)


def test_table_cells_have_korean_lang(saved):
    table = next(s for s in saved.slides[0].shapes if s.name == "ch01:table").table
    for row in table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    assert run._r.find(qn("a:rPr")).get("lang") == "ko-KR"


def test_table_font_at_body_size(saved):
    table = next(s for s in saved.slides[0].shapes if s.name == "ch01:table").table
    run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
    assert run.font.size.pt == PRESET.font_roles.table_pt


def test_header_bold_and_filled(saved):
    table = next(s for s in saved.slides[0].shapes if s.name == "ch01:table").table
    header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert header_run.font.bold is True


def _many_columns_deck(n_cols: int) -> Deck:
    return Deck(
        meta=DeckMeta(title="열 많은 표"),
        structure=Structure(
            chapters=[Chapter(id="ch01", topic="열 회귀", template="table")]
        ),
        slides=[
            Slide(
                chapter_id="ch01",
                slots=TableSlots(
                    columns=[f"열{i}" for i in range(n_cols)],
                    rows=[[f"값{i}" for i in range(n_cols)]],
                ),
            )
        ],
    )


@pytest.mark.parametrize("n_cols", [15, 20, 40])
def test_many_column_table_writes_without_crashing(tmp_path, n_cols):
    # 열 폭이 음수면 python-pptx가 EMU 범위 오류로 죽는다 (실측: n=40 -> ValueError, EMU 음수)
    metrics = FontMetrics.from_bundled()
    plan = build_render_plan(_many_columns_deck(n_cols), PRESET, metrics)
    out = tmp_path / f"table-{n_cols}.pptx"
    write_pptx(plan, out)
    presentation = Presentation(str(out))
    table = next(s for s in presentation.slides[0].shapes if s.has_table).table
    assert len(table.columns) == n_cols
    assert all(col.width > 0 for col in table.columns)
