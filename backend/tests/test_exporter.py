import json
import threading
from pathlib import Path

from pptx import Presentation

from slidecaptain.export.exporter import export_deck, export_deck_data
from slidecaptain.models.deck import (
    Bullet,
    BulletBoxSlots,
    Chapter,
    Deck,
    DeckMeta,
    Slide,
    Structure,
)
from slidecaptain.storage.file_store import FileProjectStore


def _write_deck(path: Path, title: str = "내보내기 테스트") -> Deck:
    deck = Deck(
        meta=DeckMeta(title=title),
        structure=Structure(chapters=[Chapter(id="ch01", topic="개요", template="bullet_box")]),
        slides=[
            Slide(chapter_id="ch01", slots=BulletBoxSlots(bullets=[Bullet(text="항목")], conclusion="결론"))
        ],
    )
    path.write_text(deck.model_dump_json(indent=2), encoding="utf-8")
    return deck


def test_export_creates_versioned_file(tmp_path):
    deck_path = tmp_path / "deck.json"
    _write_deck(deck_path)
    out_dir = tmp_path / "exports"
    first = export_deck(deck_path, out_dir)
    second = export_deck(deck_path, out_dir)
    assert first.name == "내보내기 테스트_v001.pptx"
    assert second.name == "내보내기 테스트_v002.pptx"
    assert first.exists() and second.exists()  # 기존 파일을 덮어쓰지 않는다


def test_export_leaves_deck_json_untouched(tmp_path):
    deck_path = tmp_path / "deck.json"
    _write_deck(deck_path)
    before = deck_path.read_bytes()
    export_deck(deck_path, tmp_path / "exports")
    assert deck_path.read_bytes() == before


def test_exported_file_opens_and_has_slides(tmp_path):
    deck_path = tmp_path / "deck.json"
    _write_deck(deck_path)
    out = export_deck(deck_path, tmp_path / "exports")
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_title_with_invalid_filename_chars_sanitized(tmp_path):
    deck_path = tmp_path / "deck.json"
    _write_deck(deck_path, title="검토: 결과/요약")
    out = export_deck(deck_path, tmp_path / "exports")
    assert out.exists()
    assert ":" not in out.name


def test_export_works_from_non_ascii_paths(tmp_path):
    korean_dir = tmp_path / "한글 폴더"
    korean_dir.mkdir()
    deck_path = korean_dir / "deck.json"
    _write_deck(deck_path, title="한글 경로 덱")
    out = export_deck(deck_path, korean_dir / "내보내기")
    assert out.exists()
    assert Presentation(str(out))


def test_bracket_title_versions_increment(tmp_path):
    # 대괄호 제목: glob 문자 클래스 해석으로 버전 스캔이 깨지던 회귀 사례
    deck_path = tmp_path / "deck.json"
    _write_deck(deck_path, title="[대외비] 검토 보고")
    out_dir = tmp_path / "exports"
    first = export_deck(deck_path, out_dir)
    second = export_deck(deck_path, out_dir)
    assert first.name.endswith("_v001.pptx")
    assert second.name.endswith("_v002.pptx")
    assert first.exists() and second.exists()


def test_preset_overrides_from_meta_applied(tmp_path):
    deck_path = tmp_path / "deck.json"
    deck = _write_deck(deck_path)
    data = json.loads(deck_path.read_text(encoding="utf-8"))
    data["meta"]["preset_overrides"] = {"font_roles": {"title_pt": 22.0}}
    deck_path.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
    out = export_deck(deck_path, tmp_path / "exports")
    prs = Presentation(str(out))
    title_shape = next(s for s in prs.slides[0].shapes if s.name == "ch01:title")
    assert title_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 22.0


def test_concurrent_export_under_store_lock_creates_all_versions(tmp_path):
    # 잠금 없이 부르면 스캔과 이동 사이의 경합으로 넷 다 v001을 돌려받는다(재현 실측).
    # store.locked(name) 안에서 부르면 저장소 잠금이 내보내기 호출을 직렬화한다 (A2가 실제 라우트에서 이렇게 부른다).
    store = FileProjectStore(tmp_path / "projects")
    store.create_project("p1", title="동시 내보내기")
    deck = _write_deck(tmp_path / "deck.json", title="동시 내보내기")
    out_dir = store.exports_dir("p1")

    results: list[Path] = []
    errors: list[Exception] = []
    results_lock = threading.Lock()

    def worker() -> None:
        try:
            with store.locked("p1"):
                path = export_deck_data(deck, out_dir)
            with results_lock:
                results.append(path)
        except Exception as e:  # noqa: BLE001 - 실패하면 아래 단언에서 드러난다
            with results_lock:
                errors.append(e)

    threads = [threading.Thread(target=worker) for _ in range(4)]
    for t in threads:
        t.start()
    for t in threads:
        t.join()

    assert errors == []
    names = sorted(p.name for p in results)
    assert names == [f"동시 내보내기_v{n:03d}.pptx" for n in range(1, 5)]
    for p in results:
        assert p.exists()
        Presentation(str(p))  # 각각 python-pptx로 정상 열린다
