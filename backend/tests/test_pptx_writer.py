import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

from slidecaptain.export.pptx_writer import write_pptx
from slidecaptain.models.preset import Preset
from slidecaptain.models.render import Frame, Para, RenderPlan, SlidePlan

PRESET = Preset()


def _simple_plan() -> RenderPlan:
    return RenderPlan(
        page_width_pt=960.0,
        page_height_pt=540.0,
        slides=[
            SlidePlan(
                chapter_id="ch01",
                template="bullet_box",
                frames=[
                    Frame(
                        name="ch01:title", x=50.0, y=36.0, w=860.0, h=40.0,
                        paras=[Para(text="장 제목", font_pt=20.0, bold=True)],
                    ),
                    Frame(
                        name="ch01:bullets", x=50.0, y=92.0, w=860.0, h=318.0,
                        paras=[
                            Para(text="첫 불릿", font_pt=12.0, bullet=True),
                            Para(text="하위 불릿", font_pt=12.0, level=1, bullet=True),
                        ],
                    ),
                    Frame(
                        name="ch01:conclusion", x=50.0, y=418.0, w=860.0, h=56.0,
                        fill="EEF3F9", border="D0D7E2",
                        paras=[Para(text="결론 문장", font_pt=12.0, bold=True, color="1F4E79")],
                    ),
                ],
            )
        ],
    )


@pytest.fixture()
def saved(tmp_path) -> Presentation:
    out = tmp_path / "out.pptx"
    write_pptx(_simple_plan(), out, PRESET)
    return Presentation(str(out))


def test_page_size_16_9(saved):
    assert saved.slide_width == Emu(12192000)
    assert saved.slide_height == Emu(6858000)


def test_shape_names_and_positions(saved):
    shapes = {s.name: s for s in saved.slides[0].shapes}
    assert set(shapes) == {"ch01:title", "ch01:bullets", "ch01:conclusion"}
    title = shapes["ch01:title"]
    # 1pt = 12700 EMU
    assert title.left == Emu(round(50.0 * 12700))
    assert title.top == Emu(round(36.0 * 12700))
    assert title.width == Emu(round(860.0 * 12700))
    assert title.height == Emu(round(40.0 * 12700))


def test_every_run_has_korean_lang_and_fonts(saved):
    for shape in saved.slides[0].shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rPr = run._r.find(qn("a:rPr"))
                assert rPr is not None
                assert rPr.get("lang") == "ko-KR"
                latin = rPr.find(qn("a:latin"))
                ea = rPr.find(qn("a:ea"))
                assert latin is not None and latin.get("typeface") == "맑은 고딕"
                assert ea is not None and ea.get("typeface") == "맑은 고딕"


def test_autofit_disabled_everywhere(saved):
    from pptx.enum.text import MSO_AUTO_SIZE

    for shape in saved.slides[0].shapes:
        if shape.has_text_frame:
            assert shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE
            assert shape.text_frame.word_wrap is True


def test_font_sizes_written_exactly(saved):
    shapes = {s.name: s for s in saved.slides[0].shapes}
    title_run = shapes["ch01:title"].text_frame.paragraphs[0].runs[0]
    assert title_run.font.size.pt == 20.0
    assert title_run.font.bold is True
    bullet_run = shapes["ch01:bullets"].text_frame.paragraphs[0].runs[0]
    assert bullet_run.font.size.pt == 12.0


def test_bullet_paragraphs_have_marker_and_level(saved):
    shapes = {s.name: s for s in saved.slides[0].shapes}
    paras = shapes["ch01:bullets"].text_frame.paragraphs
    p0 = paras[0]._p.find(qn("a:pPr"))
    assert p0 is not None
    assert p0.find(qn("a:buChar")) is not None
    assert paras[1].level == 1


def test_box_fill_and_border(saved):
    shapes = {s.name: s for s in saved.slides[0].shapes}
    box = shapes["ch01:conclusion"]
    assert box.fill.fore_color.rgb == 0xEEF3F9 or str(box.fill.fore_color.rgb) == "EEF3F9"
    assert str(box.line.color.rgb) == "D0D7E2"


def _border_only_plan() -> RenderPlan:
    return RenderPlan(
        page_width_pt=960.0,
        page_height_pt=540.0,
        slides=[
            SlidePlan(
                chapter_id="ch01",
                template="compare2",
                frames=[
                    Frame(
                        name="ch01:left_card", x=50.0, y=92.0, w=420.0, h=318.0,
                        border="D0D7E2",
                        paras=[Para(text="카드 내용", font_pt=12.0)],
                    ),
                ],
            )
        ],
    )


def test_border_only_frame_gets_padding(tmp_path):
    out = tmp_path / "border.pptx"
    write_pptx(_border_only_plan(), out, PRESET)
    prs = Presentation(str(out))
    shape = prs.slides[0].shapes[0]
    assert shape.text_frame.margin_left == Emu(round(10.0 * 12700))
