"""방법론 히스토리의 실패 목록을 고정하는 회귀 묶음 (설계서 8).

항목: 어절 줄바꿈 속성, 12pt 하한, 페이지당 크기 단계, 역할 태깅,
내보내기 전후 deck.json 불변, 결정론(같은 입력 → 같은 산출).
"""

import shutil
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from slidecaptain.export.exporter import export_deck
from slidecaptain.models.preset import BODY_MIN_PT, FOOTNOTE_MIN_PT

SAMPLE = Path(__file__).resolve().parents[1] / "samples" / "sample_deck.json"


@pytest.fixture(scope="module")
def exported(tmp_path_factory):
    work = tmp_path_factory.mktemp("regression")
    deck_path = work / "deck.json"
    shutil.copy(SAMPLE, deck_path)
    out = export_deck(deck_path, work / "exports")
    return Presentation(str(out))


def _iter_runs(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    yield from para.runs
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            yield from para.runs


def test_every_run_has_korean_lang(exported):
    runs = list(_iter_runs(exported))
    assert runs, "run이 하나도 없으면 테스트가 무의미하다"
    for run in runs:
        rPr = run._r.find(qn("a:rPr"))
        assert rPr is not None and rPr.get("lang") == "ko-KR"


def test_no_run_below_floors(exported):
    # 각주와 쪽번호는 각주 하한(9pt), 그 밖의 모든 역할은 본문 하한(12pt) 이상
    for slide in exported.slides:
        for shape in slide.shapes:
            small_ok = shape.name.endswith(":footnote") or shape.name.endswith(":page_number")
            floor = FOOTNOTE_MIN_PT if small_ok else BODY_MIN_PT
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        assert run.font.size.pt >= floor, f"{shape.name}: {run.font.size.pt}pt"
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                assert run.font.size.pt >= BODY_MIN_PT


def test_every_shape_carries_role_tag(exported):
    for slide in exported.slides:
        for shape in slide.shapes:
            assert ":" in shape.name, f"역할 태그 없는 도형: {shape.name!r}"


def test_body_area_font_steps_at_most_two_per_content_slide(exported):
    # 표지(1장)와 간지(3장)는 예외. 본문 장만 검사한다
    content_indexes = [1, 3, 4, 5]  # 0부터: summary, bullet_box, table, compare2
    for idx in content_indexes:
        slide = exported.slides[idx]
        sizes = set()
        for shape in slide.shapes:
            if shape.name.endswith(":title") or shape.name.endswith(":page_number"):
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not shape.name.endswith(":footnote"):
                            sizes.add(run.font.size.pt)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                sizes.add(run.font.size.pt)
        assert sizes, f"{idx + 1}번째 장에서 본문 run을 찾지 못했습니다"
        assert len(sizes) <= 2, f"{idx + 1}번째 장의 본문 크기 단계가 {sorted(sizes)}로 2개를 넘습니다"


def test_deterministic_export(tmp_path):
    def signature(prs):
        return [
            [
                (
                    s.name, s.left, s.top, s.width, s.height,
                    s.text_frame.text if s.has_text_frame else "",
                    tuple(cell.text_frame.text for row in s.table.rows for cell in row.cells)
                    if getattr(s, "has_table", False) and s.has_table
                    else (),
                )
                for s in slide.shapes
            ]
            for slide in prs.slides
        ]

    a_dir, b_dir = tmp_path / "a", tmp_path / "b"
    for d in (a_dir, b_dir):
        d.mkdir()
        shutil.copy(SAMPLE, d / "deck.json")
    out_a = export_deck(a_dir / "deck.json", a_dir / "exports")
    out_b = export_deck(b_dir / "deck.json", b_dir / "exports")
    assert signature(Presentation(str(out_a))) == signature(Presentation(str(out_b)))


def test_sample_deck_json_not_modified_by_export(tmp_path):
    deck_path = tmp_path / "deck.json"
    shutil.copy(SAMPLE, deck_path)
    before = deck_path.read_bytes()
    export_deck(deck_path, tmp_path / "exports")
    assert deck_path.read_bytes() == before


def test_every_text_shape_anchor_matches_render_plan_valign(exported):
    # 세로 정렬 진본은 렌더 계획이다. 채움 프레임(자동도형)이 python-pptx 기본값 ctr 을 상속해 미리보기와 어긋났던
    # 결함의 회귀 단언 (2026-09-02 Critical 묶음 태스크 B)
    from slidecaptain.layout.engine import build_render_plan
    from slidecaptain.metrics.font_metrics import FontMetrics
    from slidecaptain.models.deck import Deck
    from slidecaptain.models.preset import Preset, apply_overrides

    deck = Deck.model_validate_json(SAMPLE.read_text(encoding="utf-8"))
    plan = build_render_plan(deck, apply_overrides(Preset(), deck.meta.preset_overrides), FontMetrics.load_default())
    anchor_by_valign = {"top": "t", "middle": "ctr"}
    checked = 0
    for plan_slide, slide in zip(plan.slides, exported.slides, strict=True):
        shapes = {s.name: s for s in slide.shapes}
        for frame in plan_slide.frames:
            if frame.table is not None:
                continue
            bodyPr = shapes[frame.name].text_frame._txBody.find(qn("a:bodyPr"))
            assert (bodyPr.get("anchor") or "t") == anchor_by_valign[frame.valign], frame.name
            checked += 1
    assert checked >= 10  # 견본 덱의 텍스트 도형 전부를 실제로 검사했는지 (빈 순회 방지)
