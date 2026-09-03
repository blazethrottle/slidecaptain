"""렌더 계획 → PPTX. 축적된 기법을 내장한다 (설계서 7.1, 방법론 히스토리 C절).

- 모든 run에 ko-KR 언어 속성 (어절 단위 줄바꿈)
- a:latin과 a:ea 폰트를 함께 지정 (한글 폰트 확실 적용)
- MSO_AUTO_SIZE.NONE (자동 맞춤이 글자를 줄이는 일을 기계적으로 차단)
- 도형 이름 = 역할 태그 (향후 양방향 재수입의 열쇠)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from slidecaptain.models.render import Frame, Para, RenderPlan, RenderStyle, TablePlan

EMU_PER_PT = 12700

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
# 세로 정렬은 렌더 계획 값을 항상 명시한다. 자동도형(add_shape)의 기본값은 ctr, 텍스트박스는 top 이라
# 명시하지 않으면 채움과 테두리 프레임만 PowerPoint 에서 중앙 정렬되어 미리보기와 어긋난다 (2026-09-02 태스크 B)
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE}


def _emu(pt: float) -> Emu:
    return Emu(round(pt * EMU_PER_PT))


def _style_run(run, para: Para, style: RenderStyle) -> None:
    run.font.size = Pt(para.font_pt)
    run.font.bold = para.bold
    run.font.color.rgb = RGBColor.from_string(para.color)
    run.font.name = style.latin_font  # a:latin만 기록된다 (실측 검증 2026-08-27)
    # 공식 API가 a:rPr에 lang="ko-KR"을 기록한다 (v0.1은 한국어 고정)
    run.font.language_id = MSO_LANGUAGE_ID.KOREAN
    # 한글 폰트는 a:ea 요소로 지정해야 실제 렌더에 적용된다. 스키마 순서상 a:latin 바로 뒤에 넣는다
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        latin = rPr.find(qn("a:latin"))
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    ea.set("typeface", style.korean_font)


def _apply_bullet(paragraph, para: Para, style: RenderStyle) -> None:
    indent_emu = round(style.bullet_indent_pt * EMU_PER_PT)
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(indent_emu * (para.level + 1)))
    pPr.set("indent", str(-indent_emu))
    bu_font = pPr.makeelement(qn("a:buFont"), {"typeface": style.bullet_font})
    bu_char = pPr.makeelement(qn("a:buChar"), {"char": style.bullet_char})
    pPr.append(bu_font)
    pPr.append(bu_char)


def _fill_text_frame(tf, frame: Frame, style: RenderStyle) -> None:
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = _ANCHOR[frame.valign]
    pad = _emu(style.box_padding_pt) if (frame.fill or frame.border) else 0
    tf.margin_left = pad
    tf.margin_right = pad
    tf.margin_top = pad
    tf.margin_bottom = pad
    for i, para in enumerate(frame.paras):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        paragraph.alignment = _ALIGN[para.align]
        # 고정 pt 행간: 용량 계산(line_height_pt)과 렌더를 일치시킨다
        paragraph.line_spacing = Pt(para.font_pt * style.line_spacing)
        paragraph.level = para.level
        if para.bullet:
            _apply_bullet(paragraph, para, style)
        if i > 0 and para.bullet:
            paragraph.space_before = Pt(style.bullet_gap_pt)
        run = paragraph.add_run()
        run.text = para.text
        _style_run(run, para, style)


def _add_text_shape(slide, frame: Frame, style: RenderStyle) -> None:
    if frame.fill or frame.border:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h)
        )
        shape.shadow.inherit = False
        if frame.fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(frame.fill)
        else:
            shape.fill.background()
        if frame.border:
            shape.line.color.rgb = RGBColor.from_string(frame.border)
            shape.line.width = Pt(style.border_width_pt)
        else:
            shape.line.fill.background()
    else:
        shape = slide.shapes.add_textbox(_emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h))
    shape.name = frame.name
    _fill_text_frame(shape.text_frame, frame, style)


def _add_table_shape(slide, frame: Frame, style: RenderStyle) -> None:
    plan: TablePlan = frame.table
    n_rows = len(plan.rows) + 1
    n_cols = len(plan.header)
    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, _emu(frame.x), _emu(frame.y), _emu(frame.w), _emu(frame.h)
    )
    graphic_frame.name = frame.name
    table = graphic_frame.table
    table.first_row = False  # 내장 스타일 밴딩을 쓰지 않고 직접 칠한다 (균일성)
    table.horz_banding = False
    for i, width in enumerate(plan.col_widths_pt):
        table.columns[i].width = _emu(width)
    for i, height in enumerate(plan.row_heights_pt):
        table.rows[i].height = _emu(height)
    all_rows = [plan.header] + plan.rows
    for r_idx, row in enumerate(all_rows):
        for c_idx, text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.margin_left = _emu(style.table_cell_pad_x_pt)
            cell.margin_right = _emu(style.table_cell_pad_x_pt)
            cell.margin_top = _emu(style.table_cell_pad_y_pt)
            cell.margin_bottom = _emu(style.table_cell_pad_y_pt)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(plan.header_fill)
            tf = cell.text_frame
            tf.word_wrap = True
            paragraph = tf.paragraphs[0]
            paragraph.line_spacing = Pt(plan.font_pt * style.line_spacing)
            run = paragraph.add_run()
            run.text = text
            _style_run(
                run,
                Para(text=text, font_pt=plan.font_pt, bold=(r_idx == 0), color=style.text_color),
                style,
            )


def write_pptx(plan: RenderPlan, out_path: str | Path) -> None:
    style = plan.style
    prs = Presentation()
    prs.slide_width = _emu(plan.page_width_pt)
    prs.slide_height = _emu(plan.page_height_pt)
    blank_layout = prs.slide_layouts[6]
    for slide_plan in plan.slides:
        slide = prs.slides.add_slide(blank_layout)
        for frame in slide_plan.frames:
            if frame.table is not None:
                _add_table_shape(slide, frame, style)
            else:
                _add_text_shape(slide, frame, style)
    prs.save(str(out_path))
